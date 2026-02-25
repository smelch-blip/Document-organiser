from __future__ import annotations

import json
import os
import re
import shutil
import sys
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from rich.console import Console
from rich.progress import BarColumn, Progress, SpinnerColumn, TextColumn, TimeElapsedColumn

# Text extraction libs
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

# OpenAI
from openai import OpenAI


SUPPORTED_EXTS = {".pdf", ".pptx", ".docx", ".txt", ".md"}  # per spec  [oai_citation:5‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
MAX_CHARS_FOR_AI = 12000  # practical cap ~ "first 3,000 tokens" requirement  [oai_citation:6‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)

console = Console()


@dataclass
class AppConfig:
    input_folder: Path
    output_folder: Path
    ai_provider: str
    file_mode: str  # "move" | "copy"
    categories: List[str]


def load_config(config_path: Path) -> AppConfig:
    raw = json.loads(config_path.read_text(encoding="utf-8"))
    return AppConfig(
        input_folder=Path(raw["inputFolder"]).expanduser().resolve(),
        output_folder=Path(raw["outputFolder"]).expanduser().resolve(),
        ai_provider=str(raw.get("aiProvider", "openai")),
        file_mode=str(raw.get("fileMode", "move")).lower(),
        categories=list(raw.get("categories", ["Case Study", "Point of View", "Proposal", "Report", "Other"])),
    )


def safe_folder_name(name: str) -> str:
    # keep it simple + cross-platform safe
    name = name.strip()
    name = re.sub(r"[<>:\"/\\|?*\x00-\x1F]", "_", name)
    name = re.sub(r"\s+", " ", name)
    return name[:120] if name else "Unclassified"


def timestamp_suffix() -> str:
    # e.g. 20260215_143022 (spec example)  [oai_citation:7‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def resolve_account_folder(base: Path, account: str) -> Path:
    """
    Case-insensitive reuse of existing account folder (BR-02).  [oai_citation:8‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
    """
    desired = safe_folder_name(account)
    if not base.exists():
        return base / desired

    lower = desired.lower()
    for child in base.iterdir():
        if child.is_dir() and child.name.lower() == lower:
            return child
    return base / desired


def ensure_unique_destination(dest_path: Path) -> Path:
    if not dest_path.exists():
        return dest_path
    stem = dest_path.stem
    ext = dest_path.suffix
    return dest_path.with_name(f"{stem}_{timestamp_suffix()}{ext}")


def extract_text(file_path: Path) -> str:
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        doc = fitz.open(str(file_path))
        try:
            return "\n".join(page.get_text("text") for page in doc)
        finally:
            doc.close()

    if ext == ".pptx":
        prs = Presentation(str(file_path))
        parts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                parts.append(slide.notes_slide.notes_text_frame.text)
        return "\n".join(parts)

    if ext == ".docx":
        d = Document(str(file_path))
        return "\n".join(p.text for p in d.paragraphs if p.text)

    if ext in {".txt", ".md"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")

    return ""


def build_classifier_messages(extracted_text: str, categories: List[str]) -> Tuple[str, str]:
    """
    Prompt contract: return ONLY JSON {account,type}.  [oai_citation:9‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
    """
    sys_prompt = (
        "You are a document classification assistant. Given the text of a business document, "
        "identify: (1) the client or company name this document relates to, and (2) the document type "
        "from the list provided.\n\n"
        "Return ONLY a JSON object:\n"
        '{ "account": "<name or Unknown>", "type": "<category or Unclassified>" }\n\n'
        f"Available document types: {categories}\n\n"
        "Rules:\n"
        "- 'type' MUST be exactly one of the available document types OR 'Unclassified'.\n"
        "- If you are not confident, set account='Unknown' and type='Unclassified'.\n"
        "- Do not include any extra keys or commentary.\n"
    )
    user_msg = extracted_text[:MAX_CHARS_FOR_AI]
    return sys_prompt, user_msg


def parse_classifier_json(raw: str, categories: List[str]) -> Dict[str, str]:
    """
    If unparsable or invalid type -> Unclassified (per spec).  [oai_citation:10‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
    """
    try:
        obj = json.loads(raw.strip())
        account = str(obj.get("account", "Unknown")).strip() or "Unknown"
        doc_type = str(obj.get("type", "Unclassified")).strip() or "Unclassified"
    except Exception:
        return {"account": "Unknown", "type": "Unclassified"}

    allowed = set(categories)
    if doc_type not in allowed:
        doc_type = "Unclassified"
    if not account:
        account = "Unknown"
    return {"account": account, "type": doc_type}


def classify_with_openai(extracted_text: str, categories: List[str]) -> Dict[str, str]:
    client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
    sys_prompt, user_msg = build_classifier_messages(extracted_text, categories)

    # Using Chat Completions for widest compatibility.
    resp = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0,
        messages=[
            {"role": "system", "content": sys_prompt},
            {"role": "user", "content": user_msg},
        ],
    )

    content = resp.choices[0].message.content or ""
    return parse_classifier_json(content, categories)


def atomic_move_or_copy(src: Path, dst: Path, mode: str) -> None:
    dst.parent.mkdir(parents=True, exist_ok=True)

    if mode == "copy":
        shutil.copy2(src, dst)
        return

    # mode == "move": copy-then-delete is safer across filesystems (reliability note)  [oai_citation:11‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
    tmp_dst = dst.with_name(dst.name + ".tmp")
    if tmp_dst.exists():
        tmp_dst.unlink(missing_ok=True)
    shutil.copy2(src, tmp_dst)
    os.replace(tmp_dst, dst)
    src.unlink(missing_ok=True)


def run_sync(config: AppConfig) -> Dict:
    run_id = datetime.now(timezone.utc).isoformat()
    output_logs = config.output_folder / "logs"
    output_logs.mkdir(parents=True, exist_ok=True)  # logs folder per spec  [oai_citation:12‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)

    config.output_folder.mkdir(parents=True, exist_ok=True)

    files = [p for p in config.input_folder.iterdir() if p.is_file()]
    supported = [p for p in files if p.suffix.lower() in SUPPORTED_EXTS]
    skipped = [p for p in files if p.suffix.lower() not in SUPPORTED_EXTS]

    results: List[Dict] = []

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        BarColumn(),
        TextColumn("{task.completed}/{task.total}"),
        TimeElapsedColumn(),
        console=console,
    ) as progress:
        task = progress.add_task("Processing files", total=max(1, len(supported)))

        for idx, fpath in enumerate(supported, start=1):
            progress.update(task, description=f"({idx}/{len(supported)}) {fpath.name}")

            record = {
                "name": fpath.name,
                "account": "Unknown",
                "type": "Unclassified",
                "destination": "",
                "status": "pending",
                "error": "",
            }

            try:
                text = extract_text(fpath)
                if not text.strip():
                    classification = {"account": "Unknown", "type": "Unclassified"}
                else:
                    classification = classify_with_openai(text, config.categories)

                account = classification["account"]
                doc_type = classification["type"]

                # Unclassified fallback folder  [oai_citation:13‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
                if doc_type == "Unclassified" or account == "Unknown":
                    dest_dir = config.output_folder / "Unclassified"
                else:
                    acct_dir = resolve_account_folder(config.output_folder, account)
                    dest_dir = acct_dir / safe_folder_name(doc_type)

                dest_dir.mkdir(parents=True, exist_ok=True)  # auto-create folders  [oai_citation:14‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)

                dest_path = ensure_unique_destination(dest_dir / fpath.name)  # no overwrites  [oai_citation:15‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
                atomic_move_or_copy(fpath, dest_path, config.file_mode)

                record.update(
                    {
                        "account": account,
                        "type": doc_type,
                        "destination": str(dest_path),
                        "status": "moved" if config.file_mode == "move" else "copied",
                    }
                )
            except Exception as e:
                record.update({"status": "error", "error": str(e)})

            results.append(record)
            progress.advance(task)

    summary = {
        "runId": run_id,
        "totalFiles": len(files),
        "supportedFiles": len(supported),
        "skippedFiles": [p.name for p in skipped],
        "fileMode": config.file_mode,
        "files": results,  # log structure per spec  [oai_citation:16‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)
    }

    log_path = output_logs / f"run_{run_id.replace(':', '').replace('.', '')}.json"
    log_path.write_text(json.dumps(summary, indent=2), encoding="utf-8")

    return {"summary": summary, "logPath": str(log_path)}


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python docsort.py <path/to/config.json>")
        sys.exit(2)

    config_path = Path(sys.argv[1]).expanduser().resolve()
    if not config_path.exists():
        print(f"Config not found: {config_path}")
        sys.exit(2)

    cfg = load_config(config_path)

    cfg.input_folder.mkdir(parents=True, exist_ok=True)
    cfg.output_folder.mkdir(parents=True, exist_ok=True)

    if cfg.file_mode not in {"move", "copy"}:
        print("config.fileMode must be 'move' or 'copy'")
        sys.exit(2)

    if len(cfg.categories) < 2:
        console.print("[yellow]Warning: fewer than 2 categories defined.[/yellow]")  # per spec  [oai_citation:17‡DocSort_AI_TechSpec.docx](sediment://file_000000009df871fab687d9c3dfbb3a0d)

    if not os.environ.get("OPENAI_API_KEY"):
        console.print("[red]Missing OPENAI_API_KEY in environment.[/red]")
        sys.exit(2)

    console.print(f"[bold]Input:[/bold]  {cfg.input_folder}")
    console.print(f"[bold]Output:[/bold] {cfg.output_folder}")
    console.print(f"[bold]Mode:[/bold]   {cfg.file_mode}")
    console.print(f"[bold]Types:[/bold]  {cfg.categories}\n")

    out = run_sync(cfg)

    # Basic run summary
    files = out["summary"]["files"]
    errors = [f for f in files if f["status"] == "error"]
    unclassified = [f for f in files if f["type"] == "Unclassified" or f["account"] == "Unknown"]

    console.print("\n[bold]Run complete[/bold]")
    console.print(f"Log: {out['logPath']}")
    console.print(f"Processed: {len([f for f in files if f['status'] in {'moved','copied'}])}")
    console.print(f"Unclassified: {len(unclassified)}")
    console.print(f"Errors: {len(errors)}")


if __name__ == "__main__":
    main()
