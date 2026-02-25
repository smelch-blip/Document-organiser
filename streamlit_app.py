import json
import os
import re
import shutil
from dataclasses import dataclass, asdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import requests
import streamlit as st

# --- Optional file parsers (covered by requirements.txt) ---
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import openpyxl


APP_NAME = "DocSort AI (Streamlit)"
SUPPORTED_EXTS = {".pdf", ".pptx", ".docx", ".txt", ".md", ".xlsx", ".xlsm"}
DEFAULT_CATEGORIES = ["Case Study", "Point of View", "Proposal", "Presentation", "Report", "Other"]
MAX_CHARS_FOR_AI = 12000


# =========================
# Helpers
# =========================
def safe_name(s: str) -> str:
    s = (s or "").strip()
    s = re.sub(r"[<>:\"/\\|?*\x00-\x1F]", "_", s)
    s = re.sub(r"\s+", " ", s)
    return s[:120] if s else "Unclassified"


def ts_suffix() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


def run_id() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat() + "Z"


def ensure_unique(path: Path) -> Path:
    if not path.exists():
        return path
    return path.with_name(f"{path.stem}_{ts_suffix()}{path.suffix}")


def config_path() -> Path:
    # keep config local per-user
    base = Path.home() / ".docsort_ai"
    base.mkdir(parents=True, exist_ok=True)
    return base / "config.json"


@dataclass
class AppConfig:
    inputFolder: str = ""
    outputFolder: str = ""
    aiProvider: str = "openai"
    fileMode: str = "move"  # move or copy
    categories: List[str] = None  # type: ignore

    def __post_init__(self):
        if self.categories is None:
            self.categories = list(DEFAULT_CATEGORIES)


def load_config() -> AppConfig:
    p = config_path()
    if not p.exists():
        return AppConfig()
    try:
        data = json.loads(p.read_text(encoding="utf-8"))
        return AppConfig(
            inputFolder=data.get("inputFolder", ""),
            outputFolder=data.get("outputFolder", ""),
            aiProvider=data.get("aiProvider", "openai"),
            fileMode=data.get("fileMode", "move"),
            categories=data.get("categories", DEFAULT_CATEGORIES),
        )
    except Exception:
        return AppConfig()


def save_config(cfg: AppConfig) -> None:
    p = config_path()
    p.write_text(json.dumps(asdict(cfg), indent=2), encoding="utf-8")


# =========================
# Text extraction
# =========================
def extract_text(file_path: Path) -> str:
    ext = file_path.suffix.lower()

    try:
        if ext in {".txt", ".md"}:
            return file_path.read_text(encoding="utf-8", errors="ignore")[:MAX_CHARS_FOR_AI]

        if ext == ".docx":
            doc = Document(str(file_path))
            parts = [p.text for p in doc.paragraphs if p.text]
            return "\n".join(parts)[:MAX_CHARS_FOR_AI]

        if ext == ".pptx":
            prs = Presentation(str(file_path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    parts.append(slide.notes_slide.notes_text_frame.text)
            return "\n".join(parts)[:MAX_CHARS_FOR_AI]

        if ext == ".pdf":
            doc = fitz.open(str(file_path))
            try:
                parts = []
                for page in doc:
                    parts.append(page.get_text("text"))
                    if sum(len(x) for x in parts) >= MAX_CHARS_FOR_AI:
                        break
                return "\n".join(parts)[:MAX_CHARS_FOR_AI]
            finally:
                doc.close()

        if ext in {".xlsx", ".xlsm"}:
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            parts = []
            try:
                for ws in wb.worksheets[:5]:
                    parts.append(f"[Sheet: {ws.title}]")
                    for row in ws.iter_rows(min_row=1, max_row=50, min_col=1, max_col=12, values_only=True):
                        line = " | ".join("" if v is None else str(v) for v in row).strip()
                        if line:
                            parts.append(line)
                        if sum(len(x) for x in parts) >= MAX_CHARS_FOR_AI:
                            break
                    if sum(len(x) for x in parts) >= MAX_CHARS_FOR_AI:
                        break
                return "\n".join(parts)[:MAX_CHARS_FOR_AI]
            finally:
                try:
                    wb.close()
                except Exception:
                    pass

    except Exception:
        return ""

    return ""


# =========================
# AI Classification (OpenAI-compatible)
# =========================
def classify_openai(text: str, categories: List[str], api_key: str) -> Dict[str, str]:
    if not text.strip():
        return {"account": "Unknown", "type": "Unclassified"}

    allowed = categories or DEFAULT_CATEGORIES

    sys_prompt = (
        "You are a document classification assistant. "
        "Return ONLY a JSON object with exactly these keys:\n"
        '{ "account": "<name or Unknown>", "type": "<category or Unclassified>" }\n\n'
        f"Allowed document types: {allowed}\n\n"
        "Rules:\n"
        "- type MUST be exactly one of the allowed document types OR 'Unclassified'\n"
        "- If unsure, set account='Unknown' and type='Unclassified'\n"
        "- No extra text, no markdown."
    )

    url = os.environ.get("DOCSORT_OPENAI_URL", "https://api.openai.com/v1/chat/completions")
    model = os.environ.get("DOCSORT_OPENAI_MODEL", "gpt-4o-mini")

    payload = {
        "model": model,
        "temperature": 0.0,
        "response_format": {"type": "json_object"},
        "messages": [
            {"role": "system", "content": sys_prompt},
            {"role": "user", "content": text[:MAX_CHARS_FOR_AI]},
        ],
    }

    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}

    try:
        r = requests.post(url, headers=headers, json=payload, timeout=60)
        r.raise_for_status()
        content = r.json()["choices"][0]["message"]["content"]
        obj = json.loads(content)

        account = str(obj.get("account", "Unknown")).strip() or "Unknown"
        doc_type = str(obj.get("type", "Unclassified")).strip() or "Unclassified"

        if doc_type not in set(allowed):
            doc_type = "Unclassified"
        if not account or account.lower() in {"unknown", "n/a", "na", "none"}:
            account = "Unknown"

        return {"account": account, "type": doc_type}
    except Exception:
        return {"account": "Unknown", "type": "Unclassified"}


# =========================
# File routing + logging
# =========================
def route_destination(output_root: Path, src: Path, account: str, doc_type: str) -> Path:
    if doc_type == "Unclassified" or account == "Unknown":
        dest_dir = output_root / "Unclassified"
    else:
        dest_dir = output_root / safe_name(account) / safe_name(doc_type)

    dest_dir.mkdir(parents=True, exist_ok=True)
    return ensure_unique(dest_dir / src.name)


def transfer_file(src: Path, dst: Path, mode: str) -> str:
    try:
        dst.parent.mkdir(parents=True, exist_ok=True)
        if mode == "copy":
            shutil.copy2(src, dst)
            return "copied"
        # move (copy then delete for cross-device safety)
        shutil.copy2(src, dst)
        src.unlink(missing_ok=True)
        return "moved"
    except Exception:
        return "error"


def write_run_log(output_root: Path, run_obj: dict) -> Path:
    logs_dir = output_root / "logs"
    logs_dir.mkdir(parents=True, exist_ok=True)
    fn = f"run_{run_obj['runId'].replace(':','').replace('.','')}.json"
    p = logs_dir / fn
    p.write_text(json.dumps(run_obj, indent=2), encoding="utf-8")
    return p


# =========================
# Streamlit UI
# =========================
st.set_page_config(page_title=APP_NAME, layout="wide")
st.title(APP_NAME)

cfg = load_config()

with st.sidebar:
    st.header("Configure")

    inputFolder = st.text_input("Input folder (where files are picked up)", value=cfg.inputFolder, placeholder="e.g., C:/Temp/InputDrop")
    outputFolder = st.text_input("Output folder (where sorted files go)", value=cfg.outputFolder, placeholder="e.g., C:/Temp/SortedFiles")

    fileMode = st.selectbox("File mode", options=["move", "copy"], index=0 if cfg.fileMode == "move" else 1)
    st.caption("move = remove from input after sorting • copy = keep original")

    st.divider()

    api_key = st.text_input("OpenAI API Key", type="password", value=st.session_state.get("api_key", ""))
    st.caption("Tip: you can also set OPENAI_API_KEY in your terminal; the app will use it automatically.")

    st.divider()

    st.subheader("Categories (doc types)")
    cats_text = st.text_area(
        "One per line",
        value="\n".join(cfg.categories or DEFAULT_CATEGORIES),
        height=180,
    )
    categories = [c.strip() for c in cats_text.splitlines() if c.strip()]
    if len(categories) < 2:
        st.warning("You have fewer than 2 categories. Add more for better classification.")

    if st.button("Save settings"):
        cfg.inputFolder = inputFolder.strip()
        cfg.outputFolder = outputFolder.strip()
        cfg.fileMode = fileMode
        cfg.categories = categories
        save_config(cfg)
        st.success(f"Saved to {config_path()}")

    st.divider()
    st.subheader("Run")
    start = st.button("Sync (Classify & Sort)")

# API key resolution (session > env var)
api_key_effective = (api_key or "").strip() or os.environ.get("OPENAI_API_KEY", "").strip()

col1, col2 = st.columns([2, 1], gap="large")

with col2:
    st.subheader("What to do")
    st.markdown(
        """
1. Put files into the **Input folder**.
2. Click **Sync**.
3. Sorted files appear in the **Output folder**.
4. A JSON log is saved in `output/logs/`.
"""
    )
    st.info("This works best when you run Streamlit locally on your laptop, so it can access your folders.")

with col1:
    st.subheader("Run details")

    if start:
        inp = Path(inputFolder.strip())
        out = Path(outputFolder.strip())

        if not inp.exists() or not inp.is_dir():
            st.error("Input folder path is invalid or not a folder.")
            st.stop()
        if not out.exists() or not out.is_dir():
            st.error("Output folder path is invalid or not a folder.")
            st.stop()
        if not api_key_effective:
            st.error("Missing API key. Set it in the sidebar or via OPENAI_API_KEY environment variable.")
            st.stop()

        files = [p for p in inp.iterdir() if p.is_file()]
        supported = [p for p in files if p.suffix.lower() in SUPPORTED_EXTS]
        skipped = [p for p in files if p.suffix.lower() not in SUPPORTED_EXTS]

        st.write(f"Found **{len(files)}** files • Supported: **{len(supported)}** • Skipped: **{len(skipped)}**")
        if skipped:
            st.caption("Skipped (unsupported): " + ", ".join([p.name for p in skipped[:25]]) + (" ..." if len(skipped) > 25 else ""))

        progress = st.progress(0)
        status = st.empty()
        table = st.empty()

        results = []
        counts_by_account: Dict[str, int] = {}
        counts_by_type: Dict[str, int] = {}
        errors = []

        for i, f in enumerate(supported, start=1):
            status.write(f"Processing: **{f.name}** ({i}/{len(supported)})")
            text = extract_text(f)
            cls = classify_openai(text, categories, api_key_effective)
            account, doc_type = cls["account"], cls["type"]

            dest = route_destination(out, f, account, doc_type)
            status_str = transfer_file(f, dest, fileMode)
            if status_str == "error":
                errors.append(f.name)

            counts_by_account[account] = counts_by_account.get(account, 0) + 1
            counts_by_type[doc_type] = counts_by_type.get(doc_type, 0) + 1

            results.append(
                {
                    "file": f.name,
                    "account": account,
                    "type": doc_type,
                    "status": status_str,
                    "destination": str(dest),
                }
            )

            table.dataframe(results, use_container_width=True, hide_index=True)
            progress.progress(int(i / max(len(supported), 1) * 100))

        run_obj = {
            "runId": run_id(),
            "totalFiles": len(files),
            "supportedFiles": len(supported),
            "skippedFiles": [p.name for p in skipped],
            "fileMode": fileMode,
            "aiProvider": "openai",
            "categories": categories,
            "countsByAccount": counts_by_account,
            "countsByType": counts_by_type,
            "unclassifiedCount": sum(1 for r in results if r["type"] == "Unclassified" or r["account"] == "Unknown"),
            "errors": errors,
            "files": results,
            "endedAtLocal": datetime.now().isoformat(timespec="seconds"),
        }

        log_path = write_run_log(out, run_obj)
        status.success(f"Done. Log saved: {log_path}")

        st.download_button(
            label="Download run log JSON",
            data=json.dumps(run_obj, indent=2),
            file_name=log_path.name,
            mime="application/json",
        )
